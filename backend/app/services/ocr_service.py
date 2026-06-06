import os
import io
import logging
from azure.core.credentials import AzureKeyCredential
from azure.ai.formrecognizer import DocumentAnalysisClient

logger = logging.getLogger("somasync.ocr")

class OCRService:
    @staticmethod
    def extract_text(file_content: bytes, filename: str = "") -> str:
        """
        Extract text from file bytes (PDF, Office, Image, etc.)
        Attempts local parsing first for text-based PDF/Office files.
        Falls back to Azure Document Intelligence for scanned/image documents.
        """
        fn = filename.lower()
        extracted_text = ""
        parsed_locally = False

        try:
            # 1. Local Word document parser
            if fn.endswith((".docx", ".doc")):
                logger.info(f"Attempting local DOCX parsing for: {filename}")
                extracted_text = OCRService._parse_docx(file_content)
                parsed_locally = True

            # 2. Local PowerPoint document parser
            elif fn.endswith((".pptx", ".ppt")):
                logger.info(f"Attempting local PPTX parsing for: {filename}")
                extracted_text = OCRService._parse_pptx(file_content)
                parsed_locally = True

            # 3. Local Excel document parser
            elif fn.endswith((".xlsx", ".xls")):
                logger.info(f"Attempting local XLSX parsing for: {filename}")
                extracted_text = OCRService._parse_xlsx(file_content)
                parsed_locally = True

            # 4. Local PDF parser (with fallback to Azure if it's scanned)
            elif fn.endswith(".pdf") or file_content.startswith(b"%PDF"):
                logger.info(f"Attempting local PDF parsing for: {filename}")
                extracted_text = OCRService._parse_pdf(file_content)
                
                # Check if it looks like a scanned document (little to no text extracted)
                # We define scanned if we extract less than 100 characters of text
                cleaned_text = "".join(extracted_text.split())
                if len(cleaned_text) < 100:
                    logger.info("Local PDF extraction returned insufficient text. Falling back to Azure Document Intelligence OCR...")
                    extracted_text = "" # Trigger Azure OCR
                else:
                    parsed_locally = True

        except Exception as local_err:
            logger.warning(f"Local parser failed for {filename}: {str(local_err)}. Falling back to Azure.")
            extracted_text = ""

        # If it was parsed locally and we successfully got content, return it
        if parsed_locally and extracted_text.strip():
            return extracted_text

        # 5. Azure Document Intelligence Fallback (for images, scanned PDFs, etc.)
        logger.info(f"Running Azure Document Intelligence OCR for: {filename}")
        endpoint = os.environ.get("AZURE_DOC_INTELLIGENCE_ENDPOINT")
        key = os.environ.get("AZURE_DOC_INTELLIGENCE_KEY")

        if not endpoint or not key:
            raise ValueError(
                "Azure Document Intelligence credentials are not configured in environment. "
                "Unable to perform OCR on scanned documents/images."
            )

        # Initialize Azure Document Analysis Client
        client = DocumentAnalysisClient(
            endpoint=endpoint,
            credential=AzureKeyCredential(key)
        )

        # Use prebuilt-read for general text and OCR layout extraction
        poller = client.begin_analyze_document(
            model_id="prebuilt-read",
            document=file_content
        )
        result = poller.result()

        return result.content or ""

    @staticmethod
    def _parse_pdf(file_content: bytes) -> str:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(file_content))
        text_lines = []
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_lines.append(page_text)
        return "\n".join(text_lines)

    @staticmethod
    def _parse_docx(file_content: bytes) -> str:
        import docx
        doc = docx.Document(io.BytesIO(file_content))
        text_lines = []
        for paragraph in doc.paragraphs:
            if paragraph.text:
                text_lines.append(paragraph.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_lines.append(" | ".join(row_text))
        return "\n".join(text_lines)

    @staticmethod
    def _parse_pptx(file_content: bytes) -> str:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_content))
        text_lines = []
        for i, slide in enumerate(prs.slides):
            text_lines.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_lines.append(shape.text)
        return "\n".join(text_lines)

    @staticmethod
    def _parse_xlsx(file_content: bytes) -> str:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_content), data_only=True, read_only=True)
        text_lines = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_lines.append(f"--- Sheet: {sheet_name} ---")
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                row_count += 1
                if row_count > 100:
                    text_lines.append("[Truncated...]")
                    break
                row_vals = [str(v).strip() for v in row if v is not None]
                if row_vals:
                    text_lines.append(", ".join(row_vals))
        return "\n".join(text_lines)
